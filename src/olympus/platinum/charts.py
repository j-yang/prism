"""Chart generation helpers for PPTX slides.

Uses python-pptx native chart support for editable charts.
"""

from typing import List, Optional, Tuple

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from olympus.platinum.agent import ChartSeries, FigureSlideContent
from olympus.platinum.templates import DEFAULT_CHART_STYLE

CHART_TYPE_MAP = {
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
}


def rgb_to_rgbcolor(rgb: Tuple[int, int, int]) -> str:
    """Convert RGB tuple to RGBColor string format."""
    from pptx.dml.color import RGBColor

    return RGBColor(rgb[0], rgb[1], rgb[2])


def create_chart_data(
    x_labels: List[str],
    series_list: List[ChartSeries],
) -> CategoryChartData:
    """Create chart data object from series.

    Args:
        x_labels: X-axis category labels
        series_list: List of ChartSeries objects

    Returns:
        CategoryChartData object for pptx chart
    """
    chart_data = CategoryChartData()
    chart_data.categories = x_labels

    for series in series_list:
        chart_data.add_series(series.name, values=series.values)

    return chart_data


def add_chart_to_slide(
    slide,
    figure_content: FigureSlideContent,
    left: int,
    top: int,
    width: int,
    height: int,
    style: Optional[object] = None,
):
    """Add a native chart to a slide.

    Args:
        slide: pptx Slide object
        figure_content: FigureSlideContent with chart definition
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        style: Optional ChartStyle object
    """
    style = style or DEFAULT_CHART_STYLE

    chart_type = figure_content.chart_type.lower()
    xl_chart_type = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    if not figure_content.series:
        return None

    x_labels = (
        figure_content.x_labels
        if figure_content.x_labels
        else [str(i) for i in range(len(figure_content.series[0].values))]
    )

    chart_data = create_chart_data(x_labels, figure_content.series)

    chart = slide.shapes.add_chart(
        xl_chart_type,
        left,
        top,
        width,
        height,
        chart_data,
    ).chart

    if figure_content.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = figure_content.title
        chart.chart_title.text_frame.paragraphs[0].font.size = style.title_font_size

    if figure_content.show_legend and len(figure_content.series) > 1:
        chart.has_legend = True
        chart.legend.font.size = style.legend_font_size
    else:
        chart.has_legend = False

    if chart.value_axis:
        chart.value_axis.axis_title.text_frame.text = figure_content.y_axis_label
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = (
            style.axis_font_size
        )

    if chart.category_axis:
        chart.category_axis.axis_title.text_frame.text = figure_content.x_axis_label
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = (
            style.axis_font_size
        )

    return chart


def create_simple_bar_chart(
    slide,
    title: str,
    categories: List[str],
    values: List[float],
    series_name: str = "Values",
    left: int = 914400,
    top: int = 1600200,
    width: int = 7315200,
    height: int = 4114800,
):
    """Create a simple bar chart on a slide.

    Args:
        slide: pptx Slide object
        title: Chart title
        categories: Category labels
        values: Data values
        series_name: Series name for legend
        left, top, width, height: Position and size in EMUs

    Returns:
        Chart object
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left,
        top,
        width,
        height,
        chart_data,
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title

    return chart


def create_line_chart(
    slide,
    title: str,
    x_labels: List[str],
    series_data: List[Tuple[str, List[float]]],
    x_axis_label: str = "",
    y_axis_label: str = "",
    left: int = 914400,
    top: int = 1600200,
    width: int = 7315200,
    height: int = 4114800,
):
    """Create a line chart with multiple series.

    Args:
        slide: pptx Slide object
        title: Chart title
        x_labels: X-axis labels (time points, visits)
        series_data: List of (series_name, values) tuples
        x_axis_label: X-axis label
        y_axis_label: Y-axis label
        left, top, width, height: Position and size in EMUs

    Returns:
        Chart object
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data = CategoryChartData()
    chart_data.categories = x_labels

    for series_name, values in series_data:
        chart_data.add_series(series_name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        left,
        top,
        width,
        height,
        chart_data,
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title

    if len(series_data) > 1:
        chart.has_legend = True

    if x_axis_label and chart.category_axis:
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.text = x_axis_label

    if y_axis_label and chart.value_axis:
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = y_axis_label

    return chart
